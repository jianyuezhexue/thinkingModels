from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
import os

def create_enhanced_presentation():
    # 创建一个新的演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸为宽屏 (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 1. 标题幻灯片 - 使用背景图片
    title_slide_layout = prs.slide_layouts[6]  # 空白幻灯片
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 添加背景图片
    if os.path.exists("/Users/wang/code/thinkingModels/title_bg.png"):
        slide.shapes.add_picture("/Users/wang/code/thinkingModels/title_bg.png", 0, 0, prs.slide_width, prs.slide_height)
    
    # 添加标题文字
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(2))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "思维模型平台"
    title_p.font.size = Pt(64)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)  # 深蓝色
    title_p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9), Inches(1))
    subtitle_tf = subtitle_box.text_frame
    subtitle_p = subtitle_tf.paragraphs[0]
    subtitle_p.text = "结构化思维与AI赋能的深度思考解决方案"
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = RGBColor(102, 102, 102)  # 灰色
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # 添加标语
    slogan_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(9), Inches(1))
    slogan_tf = slogan_box.text_frame
    slogan_p = slogan_tf.paragraphs[0]
    slogan_p.text = "让思维更有结构，让思考更加深入"
    slogan_p.font.size = Pt(20)
    slogan_p.font.italic = True
    slogan_p.font.color.rgb = RGBColor(70, 70, 70)  # 深灰色
    slogan_p.alignment = PP_ALIGN.CENTER

    # 2. 项目概述幻灯片
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 添加背景
    if os.path.exists("/Users/wang/code/thinkingModels/content_bg.png"):
        slide.shapes.add_picture("/Users/wang/code/thinkingModels/content_bg.png", 0, 0, prs.slide_width, prs.slide_height)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "项目概述"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)
    
    # 添加内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    content = ("通过结构化思维和AI能力，帮助用户快速精准地完成对复杂课题的思考，\n"
               "或将想法落地。\n\n"
               "核心理念：\n"
               "• 结构化思维工具\n"
               "• AI智能引导\n"
               "• 深度思考促进\n"
               "• 解决方案导向")
    
    content_tf.text = content
    content_tf.paragraphs[0].font.size = Pt(20)
    for i in range(1, len(content_tf.paragraphs)):
        content_tf.paragraphs[i].font.size = Pt(18)
        content_tf.paragraphs[i].level = 1

    # 3. 核心功能幻灯片
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 添加背景
    if os.path.exists("/Users/wang/code/thinkingModels/content_bg.png"):
        slide.shapes.add_picture("/Users/wang/code/thinkingModels/content_bg.png", 0, 0, prs.slide_width, prs.slide_height)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "核心功能"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    content = ("1. 思维模型市场\n"
               "   • 提供精打细磨的思维模型模板\n"
               "   • 用户可上传模型并自主定价\n\n"
               
               "2. 我的课题\n"
               "   • 记录个人思考课题\n"
               "   • 选用多个思维模型辅助分析\n"
               "   • AI引导深度思考，导向解决方案\n\n"
               
               "3. 思维碰撞\n"
               "   • 公共讨论空间\n"
               "   • 思想交流与碰撞\n\n"
               
               "4. 找人聊聊\n"
               "   • 发布话题和地点\n"
               "   • 选择咖啡费用承担方式\n"
               "   • 线下面对面交流\n\n"
               
               "5. 找自己\n"
               "   • 个人思想展示\n"
               "   • 当前思考课题分享")
    
    content_tf.text = content
    content_tf.paragraphs[0].font.size = Pt(18)
    for i in range(len(content_tf.paragraphs)):
        if i == 0 or content_tf.paragraphs[i].text.strip() == "":
            content_tf.paragraphs[i].font.size = Pt(18)
        else:
            content_tf.paragraphs[i].font.size = Pt(16)
            content_tf.paragraphs[i].level = 1

    # 4. 目标用户幻灯片
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 添加背景
    if os.path.exists("/Users/wang/code/thinkingModels/content_bg.png"):
        slide.shapes.add_picture("/Users/wang/code/thinkingModels/content_bg.png", 0, 0, prs.slide_width, prs.slide_height)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "目标用户群体"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    content = ("1. 深度思考者\n"
               "   • 关注个人成长、人生规划\n"
               "   • 面临职业转折、人生成长课题\n\n"
               
               "2. 创业者及企业管理者\n"
               "   • 需要解决工作SOP问题\n"
               "   • 寻求长期问题的根本性解决方案\n"
               "   • 探索企业发展方向")
    
    content_tf.text = content
    content_tf.paragraphs[0].font.size = Pt(20)
    for i in range(1, len(content_tf.paragraphs)):
        if content_tf.paragraphs[i].text.strip() == "":
            content_tf.paragraphs[i].font.size = Pt(12)
        else:
            content_tf.paragraphs[i].font.size = Pt(18)
            content_tf.paragraphs[i].level = 1

    # 5. 商业模式幻灯片 - 带图表
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 添加背景
    if os.path.exists("/Users/wang/code/thinkingModels/content_bg.png"):
        slide.shapes.add_picture("/Users/wang/code/thinkingModels/content_bg.png", 0, 0, prs.slide_width, prs.slide_height)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "商业模式"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)
    
    # 添加图表
    if os.path.exists("/Users/wang/code/thinkingModels/chart.png"):
        slide.shapes.add_picture("/Users/wang/code/thinkingModels/chart.png", Inches(1), Inches(2), Inches(6), Inches(4.5))
    
    # 添加右侧文字说明
    desc_box = slide.shapes.add_textbox(Inches(7.5), Inches(2), Inches(5), Inches(4.5))
    desc_tf = desc_box.text_frame
    desc_tf.word_wrap = True
    
    desc = ("收入来源：\n\n"
            "• 免费用户：每周10次AI对话名额\n"
            "• 年度付费用户：89元/年\n"
            "• 终生用户：399元\n"
            "• 企业用户：20元/席位/月\n\n"
            "附加收入：\n"
            "• 模型市场抽成（10-20%）\n"
            "• 广告收入\n"
            "• 高级功能服务")
    
    desc_tf.text = desc
    for i in range(len(desc_tf.paragraphs)):
        if desc_tf.paragraphs[i].text.strip() in ["收入来源：", "附加收入："]:
            desc_tf.paragraphs[i].font.size = Pt(16)
            desc_tf.paragraphs[i].font.bold = True
        elif desc_tf.paragraphs[i].text.strip() == "":
            desc_tf.paragraphs[i].font.size = Pt(10)
        else:
            desc_tf.paragraphs[i].font.size = Pt(14)

    # 6. 技术架构幻灯片
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 添加背景
    if os.path.exists("/Users/wang/code/thinkingModels/content_bg.png"):
        slide.shapes.add_picture("/Users/wang/code/thinkingModels/content_bg.png", 0, 0, prs.slide_width, prs.slide_height)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "技术架构"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    content = ("技术栈：\n\n"
               "后端：GoLang + Gin\n"
               "   • 高性能、高并发\n"
               "   • 适合API服务\n\n"
               
               "前端：Vue + Element Plus\n"
               "   • 生态成熟\n"
               "   • 开发效率高\n\n"
               
               "数据库：MySQL + Redis\n"
               "   • 数据一致性好\n"
               "   • 读写性能优秀\n\n"
               
               "AI服务：Kim2.5 / Minimax\n"
               "   • 即插即用\n"
               "   • 无需自建模型")
    
    content_tf.text = content
    content_tf.paragraphs[0].font.size = Pt(20)
    for i in range(1, len(content_tf.paragraphs)):
        if content_tf.paragraphs[i].text.strip() in ["", "   • 高性能、高并发", "   • 生态成熟", "   • 数据一致性好", "   • 即插即用"]:
            content_tf.paragraphs[i].font.size = Pt(16)
            content_tf.paragraphs[i].level = 1
        elif content_tf.paragraphs[i].text.strip().startswith("   •"):
            content_tf.paragraphs[i].font.size = Pt(16)
            content_tf.paragraphs[i].level = 1
        else:
            content_tf.paragraphs[i].font.size = Pt(18)

    # 7. 开发路线图幻灯片
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 添加背景
    if os.path.exists("/Users/wang/code/thinkingModels/content_bg.png"):
        slide.shapes.add_picture("/Users/wang/code/thinkingModels/content_bg.png", 0, 0, prs.slide_width, prs.slide_height)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "开发路线图"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    content = ("阶段一：Web服务\n"
               "   • MVP版本开发\n"
               "   • 验证核心功能\n"
               "   • 获取早期用户反馈\n\n"
               
               "阶段二：APP应用\n"
               "   • 移动端适配\n"
               "   • 便捷随时随地思考\n\n"
               
               "阶段三：桌面端\n"
               "   • 提供强大创作功能\n"
               "   • 专业用户深度使用")
    
    content_tf.text = content
    content_tf.paragraphs[0].font.size = Pt(20)
    for i in range(1, len(content_tf.paragraphs)):
        if content_tf.paragraphs[i].text.strip() == "":
            content_tf.paragraphs[i].font.size = Pt(10)
        elif content_tf.paragraphs[i].text.strip().startswith("   •"):
            content_tf.paragraphs[i].font.size = Pt(18)
            content_tf.paragraphs[i].level = 1
        else:
            content_tf.paragraphs[i].font.size = Pt(18)

    # 8. 成本与收益分析幻灯片
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 添加背景
    if os.path.exists("/Users/wang/code/thinkingModels/content_bg.png"):
        slide.shapes.add_picture("/Users/wang/code/thinkingModels/content_bg.png", 0, 0, prs.slide_width, prs.slide_height)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "成本与收益分析"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    content = ("启动阶段（月度）：\n"
               "技术成本：约2,000元\n"
               "运营成本：约20,000元\n"
               "总成本：约22,000元\n\n"
               
               "收入预测（100用户，5%付费率）：\n"
               "月收入：约370元\n\n"
               
               "成熟阶段（月度）：\n"
               "5000用户，8%付费率：\n"
               "月收入：约29,667元\n"
               "接近盈亏平衡")
    
    content_tf.text = content
    content_tf.paragraphs[0].font.size = Pt(20)
    for i in range(1, len(content_tf.paragraphs)):
        content_tf.paragraphs[i].font.size = Pt(18)

    # 9. 结语幻灯片
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 添加背景
    if os.path.exists("/Users/wang/code/thinkingModels/title_bg.png"):
        slide.shapes.add_picture("/Users/wang/code/thinkingModels/title_bg.png", 0, 0, prs.slide_width, prs.slide_height)
    
    title_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "让思维更有结构，让思考更加深入"
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)
    title_p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9), Inches(1))
    subtitle_tf = subtitle_box.text_frame
    subtitle_p = subtitle_tf.paragraphs[0]
    subtitle_p.text = "思维模型平台 - 助力每个人成为更好的思考者"
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = RGBColor(102, 102, 102)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # 保存演示文稿
    prs.save('/Users/wang/code/thinkingModels/思维模型平台商业计划_精美版.pptx')
    print("精美版演示文稿已创建完成：/Users/wang/code/thinkingModels/思维模型平台商业计划_精美版.pptx")

if __name__ == "__main__":
    create_enhanced_presentation()